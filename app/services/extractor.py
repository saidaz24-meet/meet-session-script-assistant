from typing import Tuple, List
from pdfminer.high_level import extract_text as pdf_extract
import fitz  # PyMuPDF
from pptx import Presentation
import io

# --- PDF ---

def extract_from_pdf(file_bytes: bytes) -> str:
    try:
        return pdf_extract(io.BytesIO(file_bytes))
    except Exception:
        text = []
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        for page in doc:
            text.append(page.get_text())
        return "\n".join(text)

def pdf_page_texts(file_bytes: bytes) -> List[str]:
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    pages = []
    for page in doc:
        pages.append(page.get_text() or "")
    return pages

# --- PPTX ---

def _shape_text(shape) -> str:
    # Robust text extraction (runs/paragraphs; works better for RTL/mixed)
    if hasattr(shape, "has_text_frame") and shape.has_text_frame and shape.text_frame:
        lines = []
        for p in shape.text_frame.paragraphs:
            run_text = "".join(run.text for run in p.runs)
            if run_text.strip():
                lines.append(run_text.strip())
        return "\n".join(lines)
    if hasattr(shape, "text") and shape.text:
        return shape.text.strip()
    return ""

def extract_from_pptx(file_bytes: bytes) -> str:
    prs = Presentation(io.BytesIO(file_bytes))
    chunks = []
    for slide in prs.slides:
        slide_chunks = []
        for shape in slide.shapes:
            t = _shape_text(shape)
            if t:
                slide_chunks.append(t)
                chunks.append(t)
    return "\n".join(chunks)

def pptx_slide_texts(file_bytes: bytes) -> List[str]:
    prs = Presentation(io.BytesIO(file_bytes))
    out = []
    for slide in prs.slides:
        slide_chunks = []
        for shape in slide.shapes:
            t = _shape_text(shape)
            if t:
                slide_chunks.append(t)
        out.append("\n".join(slide_chunks))
    return out

# --- Unified ---

def extract_text(filename: str, file_bytes: bytes) -> Tuple[str, str]:
    """
    Legacy: returns (kind, joined_text)
    """
    name = filename.lower()
    if name.endswith(".pdf"):
        return "pdf", extract_from_pdf(file_bytes)
    if name.endswith(".pptx"):
        return "pptx", extract_from_pptx(file_bytes)
    try:
        return "txt", file_bytes.decode("utf-8", errors="ignore")
    except Exception:
        return "bin", ""

def extract_structured(filename: str, file_bytes: bytes) -> Tuple[str, str, List[str]]:
    """
    New: returns (kind, joined_text, per_slide_texts)
    - PDF: per_slide_texts = one entry per page
    - PPTX: per_slide_texts = one entry per slide
    - TXT/MD: single element list
    """
    name = filename.lower()
    if name.endswith(".pdf"):
        return "pdf", extract_from_pdf(file_bytes), pdf_page_texts(file_bytes)
    if name.endswith(".pptx"):
        slide_texts = pptx_slide_texts(file_bytes)
        return "pptx", "\n".join(slide_texts), slide_texts
    try:
        txt = file_bytes.decode("utf-8", errors="ignore")
        return "txt", txt, [txt]
    except Exception:
        return "bin", "", []
